import pandas as pd
import re
from pptx import Presentation
from pptx.chart.data import CategoryChartData

# === CONFIGURATION ===
EXCEL_FILE = "IMR_IBC_ Sample for Ankit- 1st July 2025(1-8) (1) 25.xlsx"
PPT_TEMPLATE = "Roadmap Delivery Confidence Assessment - TEMPLATE 93.pptx"
OUTPUT_FILE = "Roadmap Delivery Confidence Assessment.pptx"

# === LOAD DATA ===
df = pd.read_excel(EXCEL_FILE, engine='openpyxl')
total_rows = len(df)
first_row = df.iloc[0]

# === EXTRACT SOLUTION NAME FROM FILENAME ===
match = re.search(r"^(.*?)(\d{1,2}(st|nd|rd|th)?\s+\w+\s+\d{4})", EXCEL_FILE)
solution_name = match.group(1).strip().replace("_", " ") if match else "SOLUTION"

# === RELEASE INFO ===
def extract_release_parts(info):
    parts = [part.strip() for part in str(info).split(",")]
    version = parts[0] if len(parts) > 0 else "RELEASE"
    due_date = parts[1] if len(parts) > 1 else "DUE DATE"
    summary = parts[-2] if len(parts) > 2 else "PRODEL SUMMARY"
    return version, due_date, summary

release_version_1, due_date_1, prodel_summary_1 = extract_release_parts(first_row["The Release is of which version?"])
release_version_2, due_date_2, prodel_summary_2 = extract_release_parts(first_row["The Release is of which version?2"])

# === EXTRACT LAST NUMERIC VALUE FOR <KEY> ===
def extract_last_numeric(text):
    numbers = re.findall(r'\d+', str(text))
    return numbers[-1] if numbers else "KEY"

key_1 = extract_last_numeric(first_row["The Release is of which version?"])
key_2 = extract_last_numeric(first_row["The Release is of which version?2"])

# === COLUMN NAMES ===
confidence_col_1 = "How would you rate your confidence in meeting the committed release date of v1?"
lowering_col_1 = "What factors lower your confidence level?"
increasing_col_1 = "What resources, support, or actions would help increase your confidence level?"

confidence_col_2 = "How would you rate your confidence in meeting the committed release date?"
lowering_col_2 = "What factors lower your confidence level?2"
increasing_col_2 = "What resources, support, or actions would help increase your confidence level?2"

name_col = "Name"

# === COUNT CONFIDENCE RESPONSES ===
def get_confidence_counts(column):
    counts = df[column].astype(str).str.strip().str.lower().value_counts()
    return counts.get("high", 0), counts.get("medium", 0), counts.get("low", 0)

high_1, medium_1, low_1 = get_confidence_counts(confidence_col_1)
high_2, medium_2, low_2 = get_confidence_counts(confidence_col_2)

# === DETERMINE OVERALL CONFIDENCE ===
def determine_confidence(high, medium, low):
    if high == medium == low:
        return "MEDIUM"
    elif medium > high and medium > low:
        return "MEDIUM"
    elif low > high and low > medium:
        return "LOW"
    elif low == high and medium == 0:
        return "MEDIUM"
    elif low == medium and high == 0:
        return "MEDIUM"
    else:
        return "HIGH"

max_category_1 = determine_confidence(high_1, medium_1, low_1)
max_category_2 = determine_confidence(high_2, medium_2, low_2)

# === BULLET POINTS ===
def build_bullet_points(lowering_col, increasing_col):
    lowering = [f"{row.get(lowering_col, '')} — {row.get(name_col, '')}" for _, row in df.iterrows()]
    increasing = [f"{row.get(increasing_col, '')} — {row.get(name_col, '')}" for _, row in df.iterrows()]
    return lowering, increasing

lowering_points_1, increasing_points_1 = build_bullet_points(lowering_col_1, increasing_col_1)
lowering_points_2, increasing_points_2 = build_bullet_points(lowering_col_2, increasing_col_2)

# === LOAD PRESENTATION ===
prs = Presentation(PPT_TEMPLATE)

# === TEXT REPLACEMENT ===
def replace_text_in_shape(shape, replacements):
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for key, val in replacements.items():
                    if key in run.text:
                        run.text = run.text.replace(key, val)
    if shape.shape_type == 6:
        for sub_shape in shape.shapes:
            replace_text_in_shape(sub_shape, replacements)

def apply_slide_replacements(slide, replacements):
    for shape in slide.shapes:
        replace_text_in_shape(shape, replacements)

def update_chart_and_table(slide, high, medium, low):
    for shape in slide.shapes:
        if shape.has_chart:
            chart = shape.chart
            chart_data = CategoryChartData()
            chart_data.categories = ['High', 'Medium', 'Low']
            chart_data.add_series('Confidence Levels', (high, medium, low))
            chart.replace_data(chart_data)
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    val = run.text.strip().lower()
                    if val == "<high>":
                        run.text = str(high)
                    elif val == "<medium>":
                        run.text = str(medium)
                    elif val == "<low>":
                        run.text = str(low)
                    elif val in ("#", "<#>"):
                        run.text = str(total_rows)

def fill_bullet_points(slide, placeholder_text, points):
    for shape in slide.shapes:
        if shape.has_text_frame and placeholder_text in shape.text:
            shape.text_frame.clear()
            for point in points:
                shape.text_frame.add_paragraph().text = point

# === SLIDE UPDATES ===
slide1 = prs.slides[0]
apply_slide_replacements(slide1, {
    "<SOLUTION>": solution_name,
    "<Release>": release_version_1
})

slide2 = prs.slides[1]
replacements_2 = {
    "<SOLUTION>": solution_name,
    "<Release>": release_version_1,
    "<Due Date>": due_date_1,
    "<PRODEL Summary>": prodel_summary_1,
    "<Overall Confidence>": max_category_1,
    "<<VALUE>>": max_category_1,
    "<#>": str(total_rows),
    "#": str(total_rows),
    "<Key>": key_1,
    "<KEY>": key_1
}
apply_slide_replacements(slide2, replacements_2)
update_chart_and_table(slide2, high_1, medium_1, low_1)

slide3 = prs.slides[2]
apply_slide_replacements(slide3, replacements_2)
fill_bullet_points(slide3, "<<What factors lower your confidence level?>>", lowering_points_1)
fill_bullet_points(slide3, "<<What resources, support, or actions would help increase your confidence level?>>", increasing_points_1)

slide4 = prs.slides[3]
replacements_4 = {
    "<SOLUTION>": solution_name,
    "<Release>": release_version_2,
    "<Due Date>": due_date_2,
    "<PRODEL Summary>": prodel_summary_2,
    "<Overall Confidence>": max_category_2,
    "<<VALUE>>": max_category_2,
    "<#>": str(total_rows),
    "#": str(total_rows),
    "<Key>": key_2,
    "<KEY>": key_2
}
apply_slide_replacements(slide4, replacements_4)
update_chart_and_table(slide4, high_2, medium_2, low_2)

slide5 = prs.slides[4]
apply_slide_replacements(slide5, replacements_4)
fill_bullet_points(slide5, "<<What factors lower your confidence level?>>", lowering_points_2)
fill_bullet_points(slide5, "<<What resources, support, or actions would help increase your confidence level?>>", increasing_points_2)

slide6 = prs.slides[5]
apply_slide_replacements(slide6, {
    "<#>": str(total_rows),
    "<SOLUTION>": solution_name,
    "<PRODDELs>": f"{key_1},{key_2}"
})

# === SAVE PRESENTATION ===
prs.save(OUTPUT_FILE)
print(f"✅ Presentation saved as: {OUTPUT_FILE}")
